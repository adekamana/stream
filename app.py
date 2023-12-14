import collections 
import collections.abc
from pptx import Presentation
import streamlit as st
import base64
from openai import OpenAI
import pptx
from pptx.util import Pt
from bing_image_downloader import downloader
import os
from pptx.util import Inches
import time
from pathlib import Path
import json
import requests
from streamlit_lottie import st_lottie




# Change the page title
st.set_page_config(page_title='ПРЕЗЕНТАТОР', page_icon=':robot_face:')


hide_streamlit_style = """
            <style>
            #MainMenu {visibility: hidden;}
            footer {visibility: hidden;}
            </style>
            """
st.markdown(hide_streamlit_style, unsafe_allow_html=True) 










client = OpenAI(
    api_key="sk-zwZpcYNhghXEFJsbVnh4lqbLrChtawHy",
    base_url="https://api.proxyapi.ru/openai/v1",
)

def generate_slide_titles_eng(topic, num_slides):
    prompt = f"Write the title of the slides on the topic '{topic}'.The total number of slides in the presentation is '{num_slides}'. Just write the name of slides separated by commas"
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {
                "role": "user",
                "content": prompt,
            },
        ],
    )
    slide_structure = response.choices[0].message.content
    slide_titles = [part.strip() for part in slide_structure.split("\n") if part.strip()]
    return slide_titles

def generate_slide_titles(topic, num_slides):
    prompt = f"Напиши названия слайдов на тему '{topic}'. В презентации всего '{num_slides}' слайдов. Напиши только названия через запятую"
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {
                "role": "user",
                "content": prompt,
            },
        ],
    )
    slide_structure = response.choices[0].message.content
    slide_titles = [part.strip() for part in slide_structure.split("\n") if part.strip()]
    return slide_titles

def generate_slide_content(slide_title, topic):
    prompt = f"Напиши контент для слайда на тему '{topic}': '{slide_title}' До 70 слов."
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {
                "role": "user",
                "content": prompt,
            },
        ],
    )
    return response.choices[0].message.content

def generate_slide_content_eng(slide_title, topic):
    prompt = f"Write content for a slide on the topic'{topic}': '{slide_title}' Up to 70 words."
    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {
                "role": "user",
                "content": prompt,
            },
        ],
    )
    return response.choices[0].message.content

def create_presentation(topic, slide_titles, num_slides, title_font_size=32, content_font_size=17):
    prs = pptx.Presentation()

    for i in range(len(slide_titles)):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = slide_titles[i]
        title_shape.text_frame.paragraphs[0].font.size = Pt(title_font_size)
        title_shape.text_frame.paragraphs[0].font.name = "Times New Roman"

        # Generate content
        content = generate_slide_content(slide_titles[i], topic)
        
        for shape in slide.placeholders:
            if shape.name.startswith('Content Placeholder'):
                content_box = shape
                content_box.text = content[:100000]  # Limit content to 100 words
                for paragraph in content_box.text_frame.paragraphs:
                    paragraph.font.size = Pt(content_font_size)
                    paragraph.font.name = "Times New Roman"

    prs.save(f"{topic}_presentation.pptx")
    return f"{topic}_presentation.pptx"


def create_presentation_eng(topic, slide_titles, num_slides, title_font_size=32, content_font_size=17):
    prs = pptx.Presentation()

    for i in range(len(slide_titles)):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = slide_titles[i]
        title_shape.text_frame.paragraphs[0].font.size = Pt(title_font_size)
        title_shape.text_frame.paragraphs[0].font.name = "Times New Roman"

        # Generate content
        content = generate_slide_content_eng(slide_titles[i], topic)
        
        for shape in slide.placeholders:
            if shape.name.startswith('Content Placeholder'):
                content_box = shape
                content_box.text = content[:100000]  # Limit content to 100 words
                for paragraph in content_box.text_frame.paragraphs:
                    paragraph.font.size = Pt(content_font_size)
                    paragraph.font.name = "Times New Roman"

    prs.save(f"{topic}_presentation.pptx")
    return f"{topic}_presentation.pptx"


st.markdown("<h1 style='color:LightSteelBlue; font-size: 60px; text-align: center;'>ПРЕЗЕНТАТОР</h1>", unsafe_allow_html=True)



st.markdown("<h1 style='color:white; font-size: 30px; text-align: center;'>генерация презентаций Искуственным Интеллектом</h1>", unsafe_allow_html=True)






def load_lottiefile(filepath: str):
    with open(filepath, "r") as f:
        return json.load(f)

lottie_coding = load_lottiefile("lottiefile.json")  # Replace with the path to your local Lottie file

if "animation_displayed" not in st.session_state:
    st_lottie(
        lottie_coding,
        speed=1,
        reverse=False,
        loop=False,  # Set loop to False to display the animation only once
        quality="low",  # Options: "low", "medium", "high"
        height=None,
        width=None,
        key=None  # Replace with a unique key if needed
    )
    st.session_state.animation_displayed = True
    time.sleep(2)  # Display the animation for 5 seconds
    st.experimental_rerun()  # Rerun the app to remove the animation






def dynamic_text_animation(words, font_size, switch_speed, letter_speed):
    text_placeholder = st.empty()
    for word in words:
        display_word = ""
        for letter in word:
            display_word += letter
            text_placeholder.markdown(f"<h1 style='font-size:{font_size}px; text-align: center;'>{display_word}</h1>", unsafe_allow_html=True)
            time.sleep(letter_speed)  # Имитация задержки анимации ввода каждой буквы
        time.sleep(switch_speed)  # Пауза между словами
        text_placeholder.text('')  # Стирание слова


def running_dots_animation():
    placeholder = st.empty()
    while True:
        for i in range(4):
            time.sleep(0.3)
            placeholder.text("." * i)
            time.sleep(0.3)
            placeholder.text("." * (i+1))


if "has_run" not in st.session_state:
    st.session_state.has_run = True

    words = ["/Уникально/", "/Быстро/", "/Просто/"]
    font_size = 30  # Размер шрифта
    switch_speed = 1.5  # Время переключения между словами (в секундах)
    letter_speed = 0.1  # Скорость анимации ввода каждой буквы (в секундах)
    dynamic_text_animation(words, font_size, switch_speed, letter_speed)



# Use local CSS
def local_css(file_name):
    with open(file_name) as f:
        st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)


local_css("style/style.css")

# Load Animation
animation_symbol = "❄"

st.markdown(
    f"""
    <div class="snowflake">{animation_symbol}</div>
    <div class="snowflake">{animation_symbol}</div>
    <div class="snowflake">{animation_symbol}</div>
    <div class="snowflake">{animation_symbol}</div>
    <div class="snowflake">{animation_symbol}</div>
    <div class="snowflake">{animation_symbol}</div>
    <div class="snowflake">{animation_symbol}</div>
    <div class="snowflake">{animation_symbol}</div>
    <div class="snowflake">{animation_symbol}</div>
    """,
    unsafe_allow_html=True,
)





topic = st.text_input("Введите тему презентации:", max_chars=50)
num_slides = st.number_input("Введите количество слайдов:", min_value=3, max_value=19, value=6)
title_font_size = st.number_input("Укажите размер шрифта заголовка(рекоммендуется 32):", min_value=12, max_value=40, value=32)
content_font_size = st.number_input("Укажите размер шрифта текста(рекоммендуется 17):", min_value=12, max_value=32, value=17)

language_switch = st.checkbox("Включи, если хочешь сформировать презентацию на английском языке.")

generate_titles_option = st.radio("Хотите ли вы сами указать план презентации или доверите все ИИ?", ("Авто-генерация", "Ручной ввод"))

if generate_titles_option == "Авто-генерация" :
    if language_switch:
        slide_titles = generate_slide_titles_eng(topic, num_slides)
    else:
        slide_titles = generate_slide_titles(topic, num_slides)




else:
    slide_titles = []
    for i in range(num_slides):
        slide_title = st.text_input(f"Введите заголовок слайда {i+1}:")
        slide_titles.append(slide_title)

if st.button("Начать магию!") :
    if topic:
        placeholder = st.empty()
        with placeholder:
            st.write("Генерация презентации займет пару минут...")
            if language_switch:
               
                presentation_file = create_presentation_eng(topic, slide_titles, num_slides, title_font_size, content_font_size)
            else:
                
                presentation_file = create_presentation(topic, slide_titles, num_slides, title_font_size, content_font_size)        

        st.success(f"Презентация '{topic}' с {num_slides} слайдами успешно сгенерирована! Через пару секунд добавим ссылку на презентацию с изображениями!")
        
        # Add download link
        download_text = f"Скачать презентацию на тему {topic} без картинок"
        with open(presentation_file, "rb") as f:
            data = f.read()
            b64 = base64.b64encode(data).decode("utf-8")
            href = f'<a href="data:file/pptx;base64,{b64}" download="{topic}_созидание.pptx">{download_text}</a>'
            st.markdown(href, unsafe_allow_html=True)
        presentation = Presentation(f"{topic}_presentation.pptx")

        # Count the number of slides
        total_slides = len(presentation.slides)

        # Request for image search query
        query = topic

        # Download images using bing-image-downloader
        image_query = query + " slide"
        downloader.download(image_query, limit=total_slides,  output_dir='slide_images', adult_filter_off=True, force_replace=False, timeout=10)

        # Add downloaded images to the slides
        image_dir = "slide_images/"
        for i in range(total_slides):
            slide = presentation.slides[i]
            image_path = f"{image_dir}/{image_query}/image_{i+1}.jpg"
            try:
                slide.shapes.add_picture(image_path, Inches(3), Inches(4.5), width=Inches(4), height=Inches(3))
            except FileNotFoundError:
                pass  # If the image file is not found, continue to the next slide

        # Save the modified presentation
        presentation.save(f"{topic}_presentation_with_images.pptx")
        # Add download link for the presentation with images
        download_text_with_images = f"Скачать презентацию на тему {topic} с картинками"
        with open(f"{topic}_presentation_with_images.pptx", "rb") as f:
            data_with_images = f.read()
            b64_with_images = base64.b64encode(data_with_images).decode("utf-8")
            href_with_images = f'<a href="data:file/pptx;base64,{b64_with_images}" download="{topic}_созидание_с_изображениями.pptx">{download_text_with_images}</a>'
            st.markdown(href_with_images, unsafe_allow_html=True)

         

    else:
        st.error("Ошибка.Введите тему презентации!")



